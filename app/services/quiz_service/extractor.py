import os
import re
from fastapi import HTTPException


def _extract_pdf(path: str) -> str:
    import fitz  # PDF

    doc = fitz.open(path)
    pages: list[str] = []
    for page in doc:
        text = page.get_text("text")
        pages.append(text)
    doc.close()

    raw = "\n".join(pages)
    return _clean(raw)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation # powerpoint

    prs = Presentation(path)
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
    return _clean("\n".join(chunks))


def _extract_docx(path: str) -> str:
    from docx import Document # word

    doc = Document(path)
    lines = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return _clean("\n".join(lines))


def extract_text(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on server.")

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(file_path)
        else:
            raise HTTPException(
                status_code=422,
                detail=f"Quiz generation is not supported for '{ext}' files. "
                       "Please use PDF, PPTX, or DOCX.",
            )
        
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract text from the file: {exc}",
        )


# ── Pre-processing ──────────────────────────────────────────────────────────────

_PAGE_NUMBER  = re.compile(r"^\s*\d+\s*$", re.MULTILINE)
_SHORT_LINE   = re.compile(r"^.{1,3}$", re.MULTILINE)
_WHITESPACE   = re.compile(r"\n{3,}")


def _clean(text: str) -> str:
    text = _PAGE_NUMBER.sub("", text)
    text = _SHORT_LINE.sub("", text)
    text = _WHITESPACE.sub("\n\n", text)
    return text.strip()
